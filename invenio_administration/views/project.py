import os 
import requests
from langdetect import detect
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from .model import *
from sqlalchemy.sql import update
import json
# Functions
def get_matecat_languages(base_url, api_key):
    headers = {
        'accept': 'application/json',
        'x-matecat-key': api_key,
    }
    response = requests.get(f'{base_url}api/v2/languages', headers=headers)
    if response.status_code == 200:
        return response.json()
    else:
        print(f'Error fetching languages: {response.status_code}')
        return []

def find_closest_language(lang_code, matecat_languages):
    # Find the closest matching language code in MateCat API languages
    for matecat_lang in matecat_languages:
        if matecat_lang['code'].startswith(lang_code):
            return matecat_lang['code']
    return lang_code  # Return the original code if no close match is found


def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ''
    for page in reader.pages:
        text += page.extract_text() + '\n'
    return text


def detect_language(file_path, file_extension, matecat_languages):
    try:
        if file_extension == '.docx':
            text = extract_text_from_docx(file_path)
        elif file_extension == '.pptx':
            text = extract_text_from_pptx(file_path)
        elif file_extension == '.pdf':
            text = extract_text_from_pdf(file_path)
        else:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
        detected_lang = detect(text)
        mapped_lang = find_closest_language(detected_lang, matecat_languages)
        return mapped_lang
    except Exception as e:
        print(f"Error detecting language for {file_path}: {e}")
        return None

def create_project(base_url, api_key, file_path, create_data):
    headers = {
        'accept': 'application/json',
        'x-matecat-key': api_key,
    }
    with open(file_path, 'rb') as file:
        files = {'file': (file_path, file)}
        response = requests.post(f'{base_url}api/v1/new', files=files, data=create_data, headers=headers)

    if response.status_code == 200:
        print("Project is created successfully for language:", create_data['target_lang'])
        return response.json()
    else:
        print(f'Error: {response.status_code}')
        print(response.text)
        return None
    
def extract_text(element):
    text_content = element.text or ''
    for child in element:
        text_content += extract_text(child)
    return text_content

def project_creation(file_paths):
    for file_path in file_paths:
        # file_extension = os.path.splitext(file_path)[1].lower()
        file_extension="." + file_path.split('.')[-1]
        file_path_split=file_path.split('/')[-1]
        file_name_new=file_path_split.split('.')[0]
        source_lang = detect_language(file_path, file_extension, matecat_languages)
        for target_lang in target_languages:
            if target_lang[:2] == source_lang[:2]:
                continue
            create_data = {
                'project_name': 'gresis_test',
                'source_lang': source_lang,
                'target_lang': target_lang,
                'tms_engine': '1',
                'mt_engine': '1',
                'subject': 'general',
                'owner_email': 'anonymous',
                'get_public_matches': 'true',
            }

            # Create project for each target language
            project_response = create_project(base_url, api_key, file_path, create_data)
            
            if project_response['status'] == 'OK':
                uploaded = project_response['status']
                project_id = project_response['id_project']
                project_pass = project_response['project_pass']
                if(file_extension=='.json'):
                    file_name_new = file_name_new.split('_')[-1]
                original_file_row = session.query(OriginalFile).filter_by(file_name=file_name_new, project_status=None) 
                
                for original_file in original_file_row:
                    original_file_id = original_file.id
         
                if target_lang == 'ar-SA' and file_extension == '.pdf':
                    new_file_save = ArabicFile(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                elif target_lang == 'en-US' and file_extension == '.pdf':
                    new_file_save = EnglishFile(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                elif target_lang == 'fr-FR' and file_extension == '.pdf':
                    new_file_save = FrenchFile(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                elif target_lang == 'es-ES' and file_extension == '.pdf':
                    new_file_save = SpanishFile(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                elif target_lang == 'ar-SA' and file_extension == '.json':
                    new_file_save = ArabicMetadata(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                elif target_lang == 'en-US' and file_extension == '.json':
                    new_file_save = EnglishMetadata(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                elif target_lang == 'fr-FR' and file_extension == '.json':
                    new_file_save = FrenchMetadata(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                elif target_lang == 'es-ES' and file_extension == '.json':
                    new_file_save = SpanishMetadata(project_status=uploaded, project_id=project_id, project_pass=project_pass, original_file_id=original_file_id)
                session.add(new_file_save)
                session.commit()
                update_file = update(OriginalFile).where(OriginalFile.id==original_file_id).values(project_status='created')
                session.execute(update_file)
                session.commit()
                session.close()
            else:
                print(f"Failed to get project status for project ID {project_id}")
        
 
base_url = 'https://www.matecat.com/'
api_key = 'xY55uP6iraCKgrmErbQV-O4w4qe8cxXoCDnmW7Oal'
target_languages = ['ar-SA', 'en-US', 'fr-FR', 'es-ES']  

matecat_languages = get_matecat_languages(base_url, api_key)
def project_function():
    file_paths = []
    metadata_paths=[]
    files = session.query(OriginalFile).filter_by(file_type='pdf', project_status=None)  
    output_directory = "."                
    for file in files:
        dbfile_path = os.path.join(output_directory, f"{file.file_name}.pdf")
        with open(dbfile_path, "wb") as file_wr:
            file_wr.write(file.file_data)
        file_paths.append(dbfile_path)

    for metadata_file in files:
        dbfile_path = os.path.join(output_directory, f"metadata_{metadata_file.file_name}.json")
        metadata_dict = metadata_file.metadata_file
        json_data = json.dumps(metadata_dict)
        with open(dbfile_path, 'wb') as file_wr:
            file_wr.write(json_data.encode('utf-8'))
        metadata_paths.append(dbfile_path)
    project_creation(file_paths)
    project_creation(metadata_paths)
