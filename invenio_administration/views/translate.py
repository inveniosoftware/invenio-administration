import json
import requests
import re
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from .model import *

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ''
    for page in reader.pages:
        text += page.extract_text() + '\n'
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def save_text_to_file(text, file_path):
    with open(file_path, 'w', encoding='utf-8') as file:
        file.write(text)

def process_downloaded_file(url, target_language):
    response = requests.get(url)
    if response.status_code == 200:
        file = extract_file(response)
        ext = os.path.splitext(file)[1].lower()
        if ext == '.pdf':
            text = extract_text_from_pdf(file)
        elif ext == '.docx':
            text = extract_text_from_docx(file)
        elif ext == '.pptx':
            text = extract_text_from_pptx(file)
        else:
            print(f"Unsupported file format: {file}")
            return

        output_file = os.path.splitext(file)[0] + f'_{target_language}.txt'
        save_text_to_file(text, output_file)
        print(f"Extracted text from {file} to {output_file}")

def extract_file(response):
    # Check if the header is present
    cd = response.headers.get("content-disposition")
    if not cd:
        return None

    # Find filename
    fname = re.findall('filename="(.+)"', cd)
    if len(fname) == 0:
        return None
    return fname[0]

def download_file_and_save(url, target_language, project_id):
    try:
        response = requests.get(url)
        # Check if the request was successful
        if response.status_code == 200:
            filename = extract_file(response)
            if filename:
                # Save the file
                translated_filename = f"{target_language}_{filename}"
                file_name_new = translated_filename.split('.')[0]
                file_extension = translated_filename.split('.')[-1]
                # print("EXTENSION")
                # print(file_extension)
                
                # print("in file")
                with open(translated_filename, "wb") as f:
                    f.write(response.content)
                if file_extension== 'json':
                    with open(translated_filename) as fp:
                        data = json.load(fp)
                
                print(f"File downloaded and saved as {translated_filename}")
                if target_language == 'ar-SA' and file_extension == 'docx':
                    update_file = update(ArabicFile).where(ArabicFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content)
                elif target_language == 'en-US' and file_extension == 'docx':
                    update_file = update(EnglishFile).where(EnglishFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content)
                elif target_language == 'fr-FR' and file_extension == 'docx':
                    update_file = update(FrenchFile).where(FrenchFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content)
                elif target_language == 'es-ES' and file_extension == 'docx':
                    update_file = update(SpanishFile).where(SpanishFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content)
                elif target_language == 'ar-SA' and file_extension == 'json':
                    update_file = update(ArabicMetadata).where(ArabicMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data)
                elif target_language == 'en-US' and file_extension == 'json':
                    update_file = update(EnglishMetadata).where(EnglishMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data)
                elif target_language == 'fr-FR' and file_extension == 'json':
                    update_file = update(FrenchMetadata).where(FrenchMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data)
                elif target_language == 'es-ES' and file_extension == 'json':
                    update_file = update(SpanishMetadata).where(SpanishMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data)

                session.execute(update_file)
                session.commit()
                
                return translated_filename
            else:
                print("Filename could not be extracted from the response.")
        else:
            print(f"Failed to download file. HTTP status code: {response.status_code}")
    except requests.RequestException as e:
        print(f"An error occurred while downloading the file: {e}")

def extract_and_save_text(filepath):
    file_extension = os.path.splitext(filepath)[1].lower()
    extracted_text = ""
    if file_extension == '.pdf':
        extracted_text = extract_text_from_pdf(filepath)
    elif file_extension == '.docx':
        extracted_text = extract_text_from_docx(filepath)
    elif file_extension == '.pptx':
        extracted_text = extract_text_from_pptx(filepath)
    elif file_extension == '.json':
        extracted_text = ""
    else:
        print(f"Unsupported file format: {filepath}")
    
    return extracted_text

def get_download_url(base_url, api_key, job_pass, job_id):
    headers = {
        'accept': 'application/json',
        'x-matecat-key': api_key,
    }
    status_url = f'{base_url}api/v2/jobs/{job_id}/{job_pass}'
    while True:
        response = requests.get(status_url, headers=headers)
        if response.status_code == 200:
            data = response.json()
            translation_download_url = data['job']['chunks'][0]['urls']['translation_download_url']
            # print(f' retrived {translation_download_url}')
            return translation_download_url
        else:
            print(f"Failed to get job status for {job_id}: {response.status_code}")
            break

def get_job_status(base_url, api_key, status_data):
    headers = {
        'accept': 'application/json',
        'x-matecat-key': api_key,
    }
    response = requests.get(f'{base_url}api/status/', headers=headers, params=status_data)
    if response.status_code == 200:
        return response.json()
        
    else:
        print(f'Error: {response.status_code}')
        return None
    
# Main

# Get project status
def matecat_download(project_id, project_pass, uploaded):
    if uploaded == 'OK':
        # fetch downloaded status from the database
        # ..... your code here .....
        # project_id=int(project_id)
        # Demo data
        downloaded = False

        status_data = {'id_project': project_id, 'project_pass': project_pass}
        # print("STATUS DATA")
        # print(status_data)
        response_data = get_job_status(base_url, api_key, status_data)
        # print("RESPONSE")
        # print(response_data)

        # save response_data as a json file in a database
        # Demo response_data : {'errors': [], 'data': {'jobs': {'8286024': {'chunks': {'fa06541f52d0': {'12996228': {'TOTAL_PAYABLE': [368.83000000000004, '369'], 'REPETITIONS': [0, '0'], 'MT': [479, '479'], 'NEW': [0, '0'], 'TM_100': [0, '0'], 'TM_100_PUBLIC': [0, '0'], 'TM_75_99': [0, '0'], 'TM_75_84': [0, '0'], 'TM_85_94': [0, '0'], 'TM_95_99': [0, '0'], 'TM_50_74': [0, '0'], 'INTERNAL_MATCHES': [0, '0'], 'ICE': [0, '0'], 'NUMBERS_ONLY': [0, '0'], 'FILENAME': 'a.txt'}}}, 'outsource_available': True, 'totals': {'fa06541f52d0': {'TOTAL_PAYABLE': [368.83000000000004, '369'], 'REPETITIONS': [0, '0'], 'MT': [479, '479'], 'NEW': [0, '0'], 'TM_100': [0, '0'], 'TM_100_PUBLIC': [0, '0'], 'TM_75_99': [0, '0'], 'TM_75_84': [0, '0'], 'TM_85_94': [0, '0'], 'TM_95_99': [0, '0'], 'TM_50_74': [0, '0'], 'INTERNAL_MATCHES': [0, '0'], 'ICE': [0, '0'], 'NUMBERS_ONLY': [0, '0'], 'eq_word_count': [368.83000000000004, '369'], 'standard_word_count': [479, '479'], 'raw_word_count': [479, '479']}}}}, 'summary': {'IN_QUEUE_BEFORE': 0, 'STATUS': 'DONE', 'TOTAL_SEGMENTS': 31, 'SEGMENTS_ANALYZED': 31, 'TOTAL_FAST_WC': '479.00', 'TOTAL_TM_WC': 368.83000000000004, 'TOTAL_STANDARD_WC': 479, 'STANDARD_WC_TIME': '1', 'FAST_WC_TIME': '1', 'TM_WC_TIME': '59', 'STANDARD_WC_UNIT': 'hours', 'TM_WC_UNIT': 'minutes', 'FAST_WC_UNIT': 'hours', 'USAGE_FEE': '3.31', 'PRICE_PER_WORD': '0.030', 'DISCOUNT': '8', 'NAME': 'gresis_test1', 'TOTAL_RAW_WC': 479, 'TOTAL_PAYABLE': 368.83000000000004, 'PAYABLE_WC_TIME': '59', 'PAYABLE_WC_UNIT': 'minutes', 'DISCOUNT_WC': '0'}}, 'status': 'DONE', 'analyze': 'https://www.matecat.com/analyze/gresis-test1/8355507-750d4c22c675', 'jobs': {'langpairs': {'8286024-fa06541f52d0': 'en-GB|ar-SA'}, 'job-url': {'8286024-fa06541f52d0': '/translate/gresis_test1/en-GB-ar-SA/8286024-fa06541f52d0'}, 'job-quality-details': {'8286024-fa06541f52d0': [{'type': 'Typing', 'field': 'typing', 'allowed': 0.7, 'found': 0, 'founds': {'minor': '0', 'major': '0'}, 'vote': 'Excellent'}, {'type': 'Translation', 'field': 'translation', 'allowed': 0.7, 'found': 0, 'founds': {'minor': '0', 'major': '0'}, 'vote': 'Excellent'}, {'type': 'Terminology', 'field': 'terminology', 'allowed': 1.1, 'found': 0, 'founds': {'minor': '0', 'major': '0'}, 'vote': 'Excellent'}, {'type': 'Language Quality', 'field': 'language', 'allowed': 1.1, 'found': 0, 'founds': {'minor': '0', 'major': '0'}, 'vote': 'Excellent'}, {'type': 'Style', 'field': 'style', 'allowed': 1.8, 'found': 0, 'founds': {'minor': '0', 'major': '0'}, 'vote': 'Excellent'}]}, 'quality-overall': {'8286024-fa06541f52d0': 'Excellent'}}}

        # Extract project status
        job_status = response_data.get('status', 'Unknown')

        # Extract language pairs
        langpair_key = next(iter(response_data['jobs']['langpairs']), None)  # Gets the first key in the 'langpairs' dictionary
        langpair = response_data['jobs']['langpairs'].get(langpair_key) if langpair_key else 'Unknown'

        # Splitting to get target language
        if '|' in langpair:
            languages = langpair.split('|')
            target_language = languages[1] if len(languages) > 1 else 'Unknown'
        else:
            target_language = 'Unknown'
        # print("JOB STATUS")
        # print(job_status)
        if job_status == 'DONE':
            if downloaded is False:
                job_data = response_data.get('data')['jobs']
                for job_id , job_details in job_data.items():
                    job_pass = next(iter(job_details['totals']))
                    # Save job_id and job_pass in the database by creating a separate column 
                    # print(job_id)
                    # print(job_pass)
                    
                    download_translation_url = get_download_url(base_url, api_key, job_pass, job_id)
                    filename = download_file_and_save(download_translation_url, target_language, project_id)
                    file_extension= filename.split('.')[-1]
                    if file_extension != 'json':

                        searchable_text = extract_and_save_text(filename)

                        # Add this serchable_text to the invisible metadata field
                        print(searchable_text)
                
                    # update the 'downloaded' database column
                    # .... your code here ....
                    downloaded = True
            else:
                print("You have already downloaded the file")
        else:
            print(f"Error: {response_data['errors']}")
    else:
        print("Project is not uploaded yet.")


base_url = 'https://www.matecat.com/'
api_key = 'xY55uP6iraCKgrmErbQV-O4w4qe8cxXoCDnmW7Oal'

# Retrive id_project, project_pass and uploaded from database
# ..... your code here .....
def translate_function():
    arabic_records = session.query(ArabicFile).all() # Retrieve all rows from the 'files' table that are pdf
    french_records = session.query(FrenchFile).all() # Retrieve all rows from the 'files' table that are pdf
    spanish_records = session.query(SpanishFile).all()  # Retrieve all rows from the 'files' table that are pdf
    english_records = session.query(EnglishFile).all() # Retrieve all rows from the 'files' table that are pdf

    arabic_metadatas = session.query(ArabicMetadata).all() # Retrieve all rows from the 'files' table that are pdf
    french_metadatas = session.query(FrenchMetadata).all() # Retrieve all rows from the 'files' table that are pdf
    spanish_metadatas = session.query(SpanishMetadata).all()  # Retrieve all rows from the 'files' table that are pdf
    english_metadatas = session.query(EnglishMetadata).all() # Retrieve all rows from the 'files' table that are pdf


    for arabic_record in arabic_records:
        project_id = arabic_record.project_id
        project_pass = arabic_record.project_pass
        uploaded = arabic_record.project_status
        matecat_download(project_id, project_pass, uploaded)

    for french_record in french_records:
        project_id = french_record.project_id
        project_pass = french_record.project_pass
        uploaded = french_record.project_status
        matecat_download(project_id, project_pass, uploaded)

    for spanish_record in spanish_records:
        project_id = spanish_record.project_id
        project_pass = spanish_record.project_pass
        uploaded = spanish_record.project_status
        matecat_download(project_id, project_pass, uploaded)

    for english_record in english_records:
        project_id = english_record.project_id
        project_pass = english_record.project_pass
        uploaded = english_record.project_status
        matecat_download(project_id, project_pass, uploaded)

    for arabic_metadata in arabic_metadatas:
        project_id = arabic_metadata.project_id
        project_pass = arabic_metadata.project_pass
        uploaded = arabic_metadata.project_status
        matecat_download(project_id, project_pass, uploaded)

    for french_metadata in french_metadatas:
        project_id = french_metadata.project_id
        project_pass = french_metadata.project_pass
        uploaded = french_metadata.project_status
        matecat_download(project_id, project_pass, uploaded)

    for spanish_metadata in spanish_metadatas:
        project_id = spanish_metadata.project_id
        project_pass = spanish_metadata.project_pass
        uploaded = spanish_metadata.project_status
        matecat_download(project_id, project_pass, uploaded)

    for english_metadata in english_metadatas:
        project_id = english_metadata.project_id
        project_pass = english_metadata.project_pass
        uploaded = english_metadata.project_status
        matecat_download(project_id, project_pass, uploaded)

# Demo data
# project_id = 8357296
# project_pass = 'ebe81fc7b26c'
# uploaded = 'OK'

# if file_extension=='json':
                #     print('IN Json!!!')
                #     metadata_dict = response.content
                #     print(metadata_dict)
                    # translated_filename = json.dumps(metadata_dict)
                    # translated_filename = response.content
                    # with open(translated_filename, 'wb') as file_wr:
                    #     file_wr.write(metadata_dict.encode('utf-8'))
                    # bytes_data = b'This is some bytes data'

                    # Convert bytes to a string using an appropriate encoding (e.g., UTF-8)
                    # string_data = metadata_dict.decode('utf-8')

                    # # Now you can serialize the string to JSON
                    # translated_filename = json.dumps({"data": string_data})
                # else:
